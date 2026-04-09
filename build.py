#!/usr/bin/env python3
"""
Second Brain T v1.0
Transform any folder into a structured knowledge base.

What it does:
  - Scans any folder (code, docs, papers)
  - Auto-converts PDFs to markdown
  - AST parsing for Python + TypeScript/JavaScript
  - SHA256 caching (skip unchanged files)
  - Tiered context packs: Tier 0 (index), Tier 1 (topic), Tier 2 (entity)
  - Wiki with wikilinks and graph coloring (Obsidian compatible)
  - Interactive knowledge graph (vis.js) with community clustering
  - Edge confidence tagging: EXTRACTED / INFERRED / AMBIGUOUS
  - Unified dashboard linking everything
  - Auto-installs Claude skill + /sbt command

Usage:
  python3 build.py <folder>
  python3 build.py <folder> --title "My KB"
  python3 build.py <folder> --no-cache

  PDFs are auto-converted automatically — no flag needed.
"""

import os, re, sys, json, hashlib, shutil
from pathlib import Path
from datetime import datetime
from collections import defaultdict, Counter
from typing import Dict, List, Tuple, Optional

# ── Config ────────────────────────────────────────────────────────────────────

SUPPORTED_EXTS = {
    '.md', '.txt', '.rst',
    '.py', '.ts', '.tsx', '.js', '.jsx',
    '.json', '.yaml', '.yml', '.toml', '.sql', '.sh', '.css',
}
CODE_EXTS = {'.py', '.ts', '.tsx', '.js', '.jsx'}
SKIP_DIRS = {
    'node_modules', '.git', '.next', '__pycache__', 'dist', 'build',
    '.venv', 'venv', 'cache', '.obsidian', 'sbt-out', 'output',
    '.github', '.sbt-cache', 'tests',
}
SKIP_FILES = {'package-lock.json', 'yarn.lock', 'poetry.lock'}
MAX_FILE_CHARS = 10000
CACHE_DIR_NAME = '.sbt-cache'

TOPIC_COLORS = {
    'api':                '#E74C3C',
    'auth':               '#9B59B6',
    'config':             '#95A5A6',
    'database':           '#1ABC9C',
    'decision-making':    '#F5A623',
    'docs':               '#2ECC71',
    'extraction':         '#E67E22',
    'frontend':           '#3498DB',
    'graph':              '#F39C12',
    'installation':       '#27AE60',
    'megaproject-theory': '#4A90D9',
    'testing':            '#E67E22',
    'wiki':               '#8E44AD',
    'general':            '#BDC3C7',
}

STOPWORDS = {
    'the','a','an','and','or','but','in','on','at','to','for','of','with',
    'by','from','is','are','was','were','be','been','being','have','has',
    'had','do','does','did','will','would','could','should','may','might',
    'this','that','these','those','it','its','they','them','their','we',
    'our','you','your','i','my','he','she','as','if','so','not','no',
    'than','then','when','where','which','who','what','how','all','any',
    'each','can','also','more','such','into','about','up','out','one',
    'use','used','using','new','over','page','time','project','return',
    'import','class','def','function','const','let','var','type',
    'interface','export','default','async','await','true','false',
    'none','null','undefined','self','args','kwargs',
}

TOPIC_RULES = [
    (['install', 'setup', 'getting-started', 'requirements', 'install.sh'], 'installation'),
    (['api', 'route', 'endpoint', 'server', 'rest'], 'api'),
    (['auth', 'login', 'signup', 'middleware', 'session', 'oauth', 'token'], 'auth'),
    (['schema', 'database', 'db', 'sql', 'supabase', 'migration', 'model'], 'database'),
    (['component', 'ui', 'navbar', 'layout', 'page', 'frontend', 'css', 'style'], 'frontend'),
    (['config', 'tsconfig', 'package', 'next.config', 'eslint', 'pyproject', 'cargo'], 'config'),
    (['test', 'spec', 'e2e', 'fixture', 'mock'], 'testing'),
    (['extract', 'parse', 'ast', 'tree-sitter', 'ingest', 'detect', 'analyze'], 'extraction'),
    (['graph', 'cluster', 'node', 'edge', 'network', 'leiden', 'louvain', 'vis'], 'graph'),
    (['wiki', 'obsidian', 'wikilink', 'vault', 'tier', 'compiled', 'knowledge'], 'wiki'),
    (['readme', 'changelog', 'contributing', 'license', 'architecture', 'security'], 'docs'),
    (['bias', 'heuristic', 'decision', 'cognitive', 'behavior', 'behavioral'], 'decision-making'),
    (['megaproject', 'iron law', 'infrastructure', 'cost overrun', 'flyvbjerg'], 'megaproject-theory'),
]


# ── Argument parsing ──────────────────────────────────────────────────────────

def parse_args():
    args = sys.argv[1:]
    target = None
    title = 'Knowledge Base'
    use_cache = True
    i = 0
    while i < len(args):
        if args[i] == '--title' and i + 1 < len(args):
            title = args[i + 1]; i += 2
        elif args[i] == '--no-cache':
            use_cache = False; i += 1
        elif not args[i].startswith('--'):
            target = Path(args[i]).resolve(); i += 1
        else:
            i += 1
    return target, title, use_cache


# ── Helpers ───────────────────────────────────────────────────────────────────

def slugify(text: str) -> str:
    return re.sub(r'[^a-z0-9]+', '-', text.lower()).strip('-')

def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    h.update(str(path.resolve()).encode())
    try:
        h.update(path.read_bytes())
    except Exception:
        pass
    return h.hexdigest()

def strip_frontmatter(text: str) -> str:
    if text.startswith('---'):
        end = text.find('---', 3)
        if end != -1:
            return text[end + 3:].strip()
    return text

def read_file(path: Path, max_chars: int = MAX_FILE_CHARS) -> str:
    try:
        return path.read_text(encoding='utf-8', errors='ignore')[:max_chars]
    except Exception:
        return ''

def load_gitignore(folder: Path) -> List[str]:
    gi = folder / '.gitignore'
    if not gi.exists():
        return []
    return [
        line.strip().rstrip('/')
        for line in gi.read_text(errors='ignore').splitlines()
        if line.strip() and not line.startswith('#')
    ]

def is_gitignored(path: Path, root: Path, patterns: List[str]) -> bool:
    try:
        rel = str(path.relative_to(root))
    except ValueError:
        return False
    return any(p in rel or rel.endswith(p) for p in patterns)


# ── File scanning ─────────────────────────────────────────────────────────────

def collect_files(folder: Path) -> List[Path]:
    gitignore = load_gitignore(folder)
    files = []
    for root, dirs, filenames in os.walk(folder):
        root_path = Path(root)
        dirs[:] = [
            d for d in dirs
            if d not in SKIP_DIRS
            and not d.startswith('.')
            and not is_gitignored(root_path / d, folder, gitignore)
        ]
        for fname in filenames:
            if fname in SKIP_FILES:
                continue
            p = root_path / fname
            if p.suffix.lower() in SUPPORTED_EXTS and not is_gitignored(p, folder, gitignore):
                files.append(p)
    return sorted(files)

def collect_pdfs(folder: Path) -> List[Path]:
    pdfs = []
    for root, dirs, filenames in os.walk(folder):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fname in filenames:
            if fname.lower().endswith('.pdf'):
                pdfs.append(Path(root) / fname)
    return pdfs


# ── PDF conversion ────────────────────────────────────────────────────────────

def convert_pdfs(folder: Path) -> int:
    try:
        import pdfplumber
    except ImportError:
        print('  pdfplumber not installed — skipping PDFs (pip3 install pdfplumber)')
        return 0
    converted = 0
    for pdf_path in collect_pdfs(folder):
        md_path = pdf_path.with_suffix('.md')
        if md_path.exists():
            continue
        try:
            pages = []
            with pdfplumber.open(pdf_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append(f'<!-- Page {i+1} -->\n{text.strip()}')
            if pages:
                content = (f'---\ntitle: "{pdf_path.stem}"\nsource: "{pdf_path.name}"\n---\n\n'
                           f'# {pdf_path.stem}\n\n' + '\n\n---\n\n'.join(pages))
                md_path.write_text(content, encoding='utf-8')
                converted += 1
                print(f'  PDF → {md_path.name}')
        except Exception as e:
            print(f'  PDF error {pdf_path.name}: {e}')
    return converted


# ── Word conversion ───────────────────────────────────────────────────────────

def collect_docx(folder: Path) -> List[Path]:
    files = []
    for root, dirs, filenames in os.walk(folder):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fname in filenames:
            if fname.lower().endswith('.docx') and not fname.startswith('~$'):
                files.append(Path(root) / fname)
    return files

def convert_docx(folder: Path) -> int:
    try:
        from docx import Document
    except ImportError:
        print('  python-docx not installed — skipping Word docs (pip3 install python-docx)')
        return 0
    converted = 0
    for docx_path in collect_docx(folder):
        md_path = docx_path.with_suffix('.md')
        if md_path.exists():
            continue
        try:
            doc = Document(str(docx_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            if paragraphs:
                content = (f'---\ntitle: "{docx_path.stem}"\nsource: "{docx_path.name}"\n---\n\n'
                           f'# {docx_path.stem}\n\n' + '\n\n'.join(paragraphs))
                md_path.write_text(content, encoding='utf-8')
                converted += 1
                print(f'  DOCX → {md_path.name}')
        except Exception as e:
            print(f'  DOCX error {docx_path.name}: {e}')
    return converted


# ── PowerPoint conversion ─────────────────────────────────────────────────────

def collect_pptx(folder: Path) -> List[Path]:
    files = []
    for root, dirs, filenames in os.walk(folder):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fname in filenames:
            if fname.lower().endswith('.pptx') and not fname.startswith('~$'):
                files.append(Path(root) / fname)
    return files

def convert_pptx(folder: Path) -> int:
    try:
        from pptx import Presentation
    except ImportError:
        print('  python-pptx not installed — skipping PowerPoint files (pip3 install python-pptx)')
        return 0
    converted = 0
    for pptx_path in collect_pptx(folder):
        md_path = pptx_path.with_suffix('.md')
        if md_path.exists():
            continue
        try:
            prs = Presentation(str(pptx_path))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f'<!-- Slide {i} -->\n' + '\n'.join(texts))
            if slides:
                content = (f'---\ntitle: "{pptx_path.stem}"\nsource: "{pptx_path.name}"\n---\n\n'
                           f'# {pptx_path.stem}\n\n' + '\n\n---\n\n'.join(slides))
                md_path.write_text(content, encoding='utf-8')
                converted += 1
                print(f'  PPTX → {md_path.name}')
        except Exception as e:
            print(f'  PPTX error {pptx_path.name}: {e}')
    return converted


# ── SHA256 Cache ──────────────────────────────────────────────────────────────

class Cache:
    def __init__(self, root: Path, enabled: bool = True):
        self.dir = root / CACHE_DIR_NAME
        self.enabled = enabled
        if enabled:
            self.dir.mkdir(exist_ok=True)

    def load(self, path: Path) -> Optional[dict]:
        if not self.enabled:
            return None
        cache_file = self.dir / f'{file_hash(path)}.json'
        if cache_file.exists():
            try:
                return json.loads(cache_file.read_text())
            except Exception:
                return None
        return None

    def save(self, path: Path, data: dict):
        if not self.enabled:
            return
        cache_file = self.dir / f'{file_hash(path)}.json'
        tmp = cache_file.with_suffix('.tmp')
        tmp.write_text(json.dumps(data), encoding='utf-8')
        tmp.rename(cache_file)


# ── AST extraction ────────────────────────────────────────────────────────────

def extract_python_ast(path: Path) -> Tuple[List[str], List[Tuple[str, str]]]:
    try:
        import ast
        tree = ast.parse(path.read_text(encoding='utf-8', errors='ignore'))
    except Exception:
        return [], []
    nodes, edges = [], []
    label = path.stem
    for node in ast.walk(tree):
        if isinstance(node, ast.ClassDef):
            nodes.append(f'class:{node.name}')
            edges.append((label, f'class:{node.name}'))
            for base in node.bases:
                if isinstance(base, ast.Name):
                    edges.append((f'class:{node.name}', f'class:{base.id}'))
        elif isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
            nodes.append(f'fn:{node.name}')
            edges.append((label, f'fn:{node.name}'))
        elif isinstance(node, ast.Import):
            for alias in node.names:
                edges.append((label, f'import:{alias.name.split(".")[0]}'))
        elif isinstance(node, ast.ImportFrom):
            if node.module:
                edges.append((label, f'import:{node.module.split(".")[0]}'))
    return nodes, edges

def extract_ts_structure(text: str, label: str) -> List[Tuple[str, str]]:
    edges = []
    for m in re.finditer(r"import\s+.*?from\s+['\"]([^'\"]+)['\"]", text):
        mod = m.group(1)
        dep = mod.split('/')[-1].replace('.ts','').replace('.tsx','').replace('.js','')
        edges.append((label, f'{"file" if mod.startswith(".") else "pkg"}:{dep}'))
    for m in re.finditer(r'export\s+(?:default\s+)?(?:async\s+)?(?:function|class)\s+(\w+)', text):
        edges.append((label, f'export:{m.group(1)}'))
    return edges

def extract_code_structure(path: Path):
    if path.suffix == '.py':
        return extract_python_ast(path)
    elif path.suffix in {'.ts', '.tsx', '.js', '.jsx'}:
        return [], extract_ts_structure(read_file(path), path.stem)
    return [], []


# ── Topic + keyword extraction ────────────────────────────────────────────────

def infer_topic(path: Path, text: str) -> str:
    combined = ' '.join(p.lower() for p in path.parts) + ' ' + text[:500].lower()
    for keywords, topic in TOPIC_RULES:
        if any(k in combined for k in keywords):
            return topic
    return 'general'

def extract_keywords(text: str, top_n: int = 15) -> List[str]:
    words = re.sub(r'[^a-zA-Z\s]', ' ', text.lower()).split()
    freq = Counter(w for w in words if len(w) > 4 and w not in STOPWORDS)
    return [w for w, _ in freq.most_common(top_n)]


# ── Entity building ───────────────────────────────────────────────────────────

def build_entities(files: List[Path], target: Path, cache: Cache) -> List[dict]:
    entities = []
    for f in files:
        cached = cache.load(f)
        if cached:
            entities.append(cached)
            continue
        text = read_file(f)
        if not text.strip():
            continue
        clean = strip_frontmatter(text)
        ast_nodes, ast_edges = [], []
        if f.suffix in CODE_EXTS:
            ast_nodes, ast_edges = extract_code_structure(f)
        entity = {
            'id': slugify(f.stem),
            'label': f.stem.replace('-', ' ').replace('_', ' ').title(),
            'file': str(f.relative_to(target)),
            'ext': f.suffix,
            'text': clean,
            'keywords': extract_keywords(clean),
            'topic': infer_topic(f, clean),
            'ast_nodes': ast_nodes,
            'ast_edges': [[a, b] for a, b in ast_edges],
            'is_code': f.suffix in CODE_EXTS,
        }
        cache.save(f, entity)
        entities.append(entity)
    # deduplicate IDs
    seen: Dict[str, int] = {}
    for e in entities:
        orig = e['id']
        n = seen.get(orig, 0)
        if n > 0:
            e['id'] = f'{orig}-{n}'
        seen[orig] = n + 1
    return entities


# ── Relationship detection ────────────────────────────────────────────────────

def find_relationships(entities: List[dict]) -> List[dict]:
    id_map = {e['id']: e for e in entities}
    edges = []
    seen_pairs = set()

    # AST edges — EXTRACTED
    for e in entities:
        for src, tgt in e.get('ast_edges', []):
            tgt_slug = slugify(tgt.split(':')[-1])
            if tgt_slug in id_map and tgt_slug != e['id']:
                pair = tuple(sorted([e['id'], tgt_slug]))
                if pair not in seen_pairs:
                    seen_pairs.add(pair)
                    edges.append({
                        'from': e['id'], 'to': tgt_slug,
                        'weight': 10, 'type': 'EXTRACTED',
                        'shared_keywords': [], 'relation': 'imports',
                    })

    # Semantic edges
    elist = list(entities)
    for i, a in enumerate(elist):
        for j, b in enumerate(elist):
            if j <= i:
                continue
            pair = tuple(sorted([a['id'], b['id']]))
            shared = list(set(a['keywords']) & set(b['keywords']))
            same_topic = a['topic'] == b['topic']
            a_txt = a['text'][:3000].lower()
            b_txt = b['text'][:3000].lower()
            a_mentions_b = len(b['label']) > 3 and b['label'].lower().split()[0] in a_txt
            b_mentions_a = len(a['label']) > 3 and a['label'].lower().split()[0] in b_txt
            weight = (
                len(shared) * 2 +
                (3 if same_topic else 0) +
                (5 if a_mentions_b else 0) +
                (5 if b_mentions_a else 0)
            )
            if weight >= 4:
                edge_type = ('EXTRACTED' if pair in seen_pairs
                             else 'INFERRED' if weight >= 8
                             else 'AMBIGUOUS')
                if pair not in seen_pairs:
                    seen_pairs.add(pair)
                    edges.append({
                        'from': a['id'], 'to': b['id'],
                        'weight': weight, 'type': edge_type,
                        'shared_keywords': shared[:5],
                        'relation': 'same-topic' if same_topic else 'keyword-overlap',
                    })
    return edges


# ── Community detection (label propagation, no external deps) ─────────────────

def detect_communities(entities: List[dict], edges: List[dict]) -> Dict[str, int]:
    import random
    ids = [e['id'] for e in entities]
    adj: Dict[str, set] = defaultdict(set)
    for edge in edges:
        if edge['weight'] >= 6:
            adj[edge['from']].add(edge['to'])
            adj[edge['to']].add(edge['from'])
    labels = {n: i for i, n in enumerate(ids)}
    for _ in range(25):
        changed = False
        order = ids[:]
        random.shuffle(order)
        for node in order:
            nbrs = list(adj[node])
            if not nbrs:
                continue
            nl = [labels[nb] for nb in nbrs]
            best = max(set(nl), key=nl.count)
            if labels[node] != best:
                labels[node] = best
                changed = True
        if not changed:
            break
    counts = Counter(labels.values())
    rank = {label: i for i, (label, _) in enumerate(counts.most_common())}
    return {node: rank[label] for node, label in labels.items()}


# ── Tier generation ──────────────────────────────────────────────────────────────

def build_tiers(entities: List[dict], edges: List[dict], out_dir: Path, title: str):
    out_dir.mkdir(parents=True, exist_ok=True)
    topics: Dict[str, list] = defaultdict(list)
    for e in entities:
        topics[e['topic']].append(e)

    # Tier 0 — global index
    lines = [
        f'# {title} — Tier 0: Index\n',
        f'_Compiled: {datetime.now().strftime("%Y-%m-%d %H:%M")} | Entities: {len(entities)} | Connections: {len(edges)}_\n\n',
        '## All Entities\n',
    ]
    for e in entities:
        lines.append(f'- **{e["label"]}** ({e["topic"]}) — {", ".join(e["keywords"][:3])}')
    lines.append('\n## Topics\n')
    for topic, ents in sorted(topics.items()):
        lines.append(f'- **{topic}** ({len(ents)} entities)')
    (out_dir / 'index.md').write_text('\n'.join(lines), encoding='utf-8')

    # Tier 1 — by topic
    t1 = out_dir / 'topic'
    t1.mkdir(exist_ok=True)
    for topic, ents in topics.items():
        lines = [f'# {topic.replace("-"," ").title()} — Tier 1: Topic Context\n', f'_Entities: {len(ents)}_\n']
        for e in ents:
            content = e['text'].strip()[:2000]
            if len(e['text'].strip()) > 2000:
                content += '\n\n_[truncated — see Tier 2]_'
            lines += [f'\n## {e["label"]} (`{e["file"]}`)\n',
                      f'_Keywords: {", ".join(e["keywords"][:5])}_\n', content]
        (t1 / f'{slugify(topic)}.md').write_text('\n'.join(lines), encoding='utf-8')

    # Tier 2 — per entity
    t2 = out_dir / 'entity'
    t2.mkdir(exist_ok=True)
    edge_map: Dict[str, list] = defaultdict(list)
    for edge in edges:
        edge_map[edge['from']].append(edge)
        edge_map[edge['to']].append(edge)
    id_map = {e['id']: e for e in entities}
    for e in entities:
        related_edges = sorted(edge_map[e['id']], key=lambda x: -x['weight'])[:8]
        related = [(id_map[ed['to'] if ed['from'] == e['id'] else ed['from']], ed)
                   for ed in related_edges if (ed['to'] if ed['from'] == e['id'] else ed['from']) in id_map]
        lines = [
            f'# {e["label"]} — Tier 2: Deep Context\n',
            f'_File: `{e["file"]}` | Topic: {e["topic"]}_\n',
            f'_Keywords: {", ".join(e["keywords"][:8])}_\n',
        ]
        if e.get('ast_nodes'):
            lines += ['\n## Code Structure\n'] + [f'- `{n}`' for n in e['ast_nodes'][:12]]
        lines += ['\n## Content\n', e['text'].strip()[:5000] + ('\n\n_[truncated]_' if len(e['text']) > 5000 else '')]
        if related:
            lines.append('\n## Connections\n')
            for other, ed in related:
                lines.append(f'- **{other["label"]}** [{ed["type"]}] — {", ".join(ed["shared_keywords"][:3]) or ed["relation"]}')
        (t2 / f'{e["id"]}.md').write_text('\n'.join(lines), encoding='utf-8')

    print(f'  Tiers: index + {len(topics)} topic files + {len(entities)} entity files')


# ── Wiki generation ───────────────────────────────────────────────────────────

def build_wiki(entities: List[dict], edges: List[dict], wiki_dir: Path, title: str):
    wiki_dir.mkdir(parents=True, exist_ok=True)
    edge_map: Dict[str, list] = defaultdict(list)
    for edge in edges:
        edge_map[edge['from']].append(edge)
        edge_map[edge['to']].append(edge)
    id_map = {e['id']: e for e in entities}
    topics: Dict[str, list] = defaultdict(list)
    for e in entities:
        topics[e['topic']].append(e)

    for e in entities:
        topic = e['topic']
        (wiki_dir / topic).mkdir(exist_ok=True)
        related_edges = sorted(edge_map[e['id']], key=lambda x: -x['weight'])[:6]
        related = [(id_map[ed['to'] if ed['from'] == e['id'] else ed['from']], ed)
                   for ed in related_edges if (ed['to'] if ed['from'] == e['id'] else ed['from']) in id_map]
        content = re.sub(r'\n{3,}', '\n\n', e['text'].strip())
        if len(content) > 4000:
            content = content[:4000] + '\n\n_[truncated]_'
        lines = [
            f'---\ntitle: "{e["label"]}"\ntopic: {topic}\ntags: [{topic}]\nfile: {e["file"]}\n---\n\n',
            f'# {e["label"]}\n\n',
            f'> **File:** `{e["file"]}` | **Topic:** {topic}\n\n',
        ]
        if e.get('ast_nodes'):
            lines += ['## Structure\n'] + [f'- `{n}`' for n in e['ast_nodes'][:8]] + ['\n']
        lines.append(content)
        if related:
            lines.append('\n\n## Related\n')
            for other, ed in related:
                tag = f' [{ed["type"]}]' if ed['type'] == 'EXTRACTED' else ''
                shared = ', '.join(ed['shared_keywords'][:3])
                lines.append(f'- [[{other["topic"]}/{other["id"]}|{other["label"]}]]{tag} — _{shared or ed["relation"]}_')
        lines.append(f'\n\n---\n_[View in graph](../../graph/graph.html) | [Deep context](../../tiers/entity/{e["id"]}.md)_\n')
        (wiki_dir / topic / f'{e["id"]}.md').write_text(''.join(lines), encoding='utf-8')

    for topic, ents in topics.items():
        lines = [f'---\ntags: [{topic}, overview]\n---\n\n# {topic.replace("-"," ").title()}\n\n']
        for e in ents:
            lines.append(f'- [[{topic}/{e["id"]}|{e["label"]}]] ({len(edge_map[e["id"]])} connections)')
        (wiki_dir / topic / '_overview.md').write_text('\n'.join(lines), encoding='utf-8')

    moc = [f'---\ntags: [index]\n---\n\n# {title} — Map of Content\n\n',
           f'_Generated: {datetime.now().strftime("%Y-%m-%d %H:%M")}_\n']
    for topic, ents in sorted(topics.items()):
        moc.append(f'\n## {topic.replace("-"," ").title()} ({len(ents)})\n')
        for e in ents:
            moc.append(f'- [[{topic}/{e["id"]}|{e["label"]}]]')
    (wiki_dir / 'index.md').write_text('\n'.join(moc), encoding='utf-8')

    obs = wiki_dir / '.obsidian'
    obs.mkdir(exist_ok=True)
    groups = [{'query': f'path:{t}', 'color': {'a': 1, 'rgb': int(TOPIC_COLORS.get(t, '#888888')[1:], 16)}, 'lineStyle': 0}
              for t in topics if t in TOPIC_COLORS]
    (obs / 'graph.json').write_text(json.dumps({'colorGroups': groups, 'showTags': True}, indent=2))
    (obs / 'app.json').write_text(json.dumps({'showLineNumber': True, 'readableLineLength': True, 'showFrontmatter': False}, indent=2))

    total = len(list(wiki_dir.rglob('*.md')))
    print(f'  Wiki: {total} pages across {len(topics)} topics')


# ── Graph (vis.js + report) ───────────────────────────────────────────────────

def build_graph(entities: List[dict], edges: List[dict], graph_dir: Path,
                title: str, communities: Dict[str, int]):
    graph_dir.mkdir(parents=True, exist_ok=True)
    edge_map: Dict[str, list] = defaultdict(list)
    for edge in edges:
        edge_map[edge['from']].append(edge)
        edge_map[edge['to']].append(edge)

    nodes = [{
        'id': e['id'], 'label': e['label'], 'topic': e['topic'], 'file': e['file'],
        'keywords': e['keywords'][:8], 'degree': len(edge_map[e['id']]),
        'community': communities.get(e['id'], 0),
        'color': TOPIC_COLORS.get(e['topic'], '#BDC3C7'),
        'size': 12 + min(len(edge_map[e['id']]) * 3, 28),
        'is_code': e.get('is_code', False),
    } for e in entities]

    n_comm = max(communities.values(), default=0) + 1
    stats = {
        'total_nodes': len(nodes), 'total_edges': len(edges), 'communities': n_comm,
        'topics': sorted(set(e['topic'] for e in entities)),
        'extracted': sum(1 for e in edges if e['type'] == 'EXTRACTED'),
        'inferred': sum(1 for e in edges if e['type'] == 'INFERRED'),
        'ambiguous': sum(1 for e in edges if e['type'] == 'AMBIGUOUS'),
    }
    (graph_dir / 'graph.json').write_text(json.dumps(
        {'generated_at': datetime.now().isoformat(), 'title': title,
         'nodes': nodes, 'edges': edges, 'stats': stats}, indent=2), encoding='utf-8')

    # Report
    id_to_label = {n['id']: n['label'] for n in nodes}
    hub_nodes = sorted(nodes, key=lambda n: -n['degree'])[:5]
    strong_edges = sorted([e for e in edges if e['weight'] >= 8], key=lambda x: -x['weight'])[:5]
    topic_counts = Counter(e['topic'] for e in entities)
    isolated = [n for n in nodes if n['degree'] == 0]

    report = [
        f'# {title} — Graph Report\n',
        f'_Generated: {datetime.now().strftime("%Y-%m-%d %H:%M")}_\n\n',
        f'## Statistics\n',
        f'- **Nodes:** {stats["total_nodes"]} | **Connections:** {stats["total_edges"]} | **Clusters:** {stats["communities"]}',
        f'- **Edge confidence:** EXTRACTED (AST): {stats["extracted"]} | INFERRED: {stats["inferred"]} | AMBIGUOUS: {stats["ambiguous"]}\n',
        f'## Hub Nodes (most connected)\n',
    ] + [f'- **{n["label"]}** ({n["topic"]}) — {n["degree"]} connections' for n in hub_nodes]
    report += [f'\n## Topics\n'] + [f'- **{t.replace("-"," ").title()}**: {c} entities' for t, c in topic_counts.most_common()]
    if strong_edges:
        report += [f'\n## Strongest Connections\n'] + [
            f'- **{id_to_label.get(e["from"], e["from"])}** ↔ **{id_to_label.get(e["to"], e["to"])}** [{e["type"]}] — {", ".join(e["shared_keywords"][:3]) or e["relation"]}'
            for e in strong_edges]
    if isolated:
        report += [f'\n## Knowledge Gaps ({len(isolated)} isolated nodes)\n'] + [
            f'- **{n["label"]}** has no connections' for n in isolated[:5]]
    report += [
        f'\n## Suggested Questions\n',
        f'- What is the role of **{hub_nodes[0]["label"]}** in this project?' if hub_nodes else '',
        f'- How does **{hub_nodes[0]["label"]}** connect to **{hub_nodes[1]["label"]}**?' if len(hub_nodes) > 1 else '',
        f'- What are the key themes across all {len(stats["topics"])} topics?',
        f'- Review {stats["ambiguous"]} ambiguous connections for accuracy' if stats['ambiguous'] else '',
        f'\n---\n_[Open graph](graph.html) | [Wiki](../wiki/index.md) | [Tiers](../tiers/index.md)_\n',
    ]
    (graph_dir / 'report.md').write_text('\n'.join(r for r in report if r is not None), encoding='utf-8')
    print(f'  Report: graph/report.md')

    # HTML
    topics_list = sorted(set(n['topic'] for n in nodes))
    vis_nodes = json.dumps([{
        'id': n['id'], 'label': n['label'][:28],
        'title': (f"<b>{n['label']}</b><br>Topic: {n['topic']}<br>File: {n['file']}<br>"
                  f"Connections: {n['degree']}<br>Cluster: {n['community']}<br>"
                  f"Keywords: {', '.join(n['keywords'][:5])}"),
        'color': {'background': n['color'], 'border': '#ffffff33',
                  'highlight': {'background': n['color'], 'border': '#fff'}},
        'size': n['size'], 'font': {'color': '#fff', 'size': 12},
        'topic': n['topic'], 'community': n['community'], 'is_code': n['is_code'],
    } for n in nodes])
    vis_edges = json.dumps([{
        'from': e['from'], 'to': e['to'],
        'width': min(1 + e['weight'] / 5, 5),
        'title': f"{e['type']}: {', '.join(e['shared_keywords'][:3]) or e['relation']}",
        'color': {'color': '#388bfd' if e['type'] == 'EXTRACTED' else ('#8b949e' if e['type'] == 'INFERRED' else '#484f58'),
                  'highlight': '#e6edf3'},
        'dashes': e['type'] == 'AMBIGUOUS',
    } for e in edges])
    topic_filters = '\n'.join(
        f'<label><input type="checkbox" class="ft" value="{t}" checked>'
        f'<span style="color:{TOPIC_COLORS.get(t,"#888")}">■</span> {t.replace("-"," ").title()}</label>'
        for t in topics_list)
    comm_filters = '\n'.join(
        f'<label><input type="checkbox" class="fc" value="{i}" checked> Cluster {i}</label>'
        for i in range(n_comm))

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><title>{title} — Graph</title>
<script src="https://unpkg.com/vis-network/standalone/umd/vis-network.min.js"></script>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;background:#0d1117;color:#e6edf3;height:100vh;display:flex;flex-direction:column}}
#hdr{{padding:12px 20px;background:#161b22;border-bottom:1px solid #30363d;display:flex;align-items:center;justify-content:space-between;flex-shrink:0}}
#hdr h1{{font-size:16px;font-weight:600}}
.st{{font-size:12px;color:#8b949e;display:flex;gap:14px;align-items:center}}
.badge{{padding:2px 8px;border-radius:10px;font-size:11px;font-weight:600}}
.EX{{background:#1f6feb33;color:#58a6ff}}.IN{{background:#388bfd22;color:#79c0ff}}.AM{{background:#f0883e22;color:#f0883e}}
#body{{display:flex;flex:1;overflow:hidden}}
#sb{{width:210px;background:#161b22;border-right:1px solid #30363d;overflow-y:auto;padding:12px;flex-shrink:0;font-size:12px}}
#sb h4{{font-size:10px;text-transform:uppercase;letter-spacing:.08em;color:#8b949e;margin:10px 0 5px}}
#sb label{{display:flex;align-items:center;gap:6px;margin-bottom:4px;cursor:pointer;color:#c9d1d9}}
#search{{width:100%;padding:6px 9px;background:#0d1117;border:1px solid #30363d;border-radius:6px;color:#e6edf3;font-size:12px;margin-bottom:6px}}
#search:focus{{outline:none;border-color:#388bfd}}
#graph{{flex:1}}
#info{{width:240px;background:#161b22;border-left:1px solid #30363d;padding:14px;overflow-y:auto;flex-shrink:0;font-size:12px;display:none}}
#info h3{{font-size:14px;font-weight:600;color:#fff;margin-bottom:6px;word-break:break-word}}
.meta{{color:#8b949e;line-height:1.7;margin-bottom:8px}}
.kws span{{display:inline-block;background:#21262d;padding:2px 7px;border-radius:10px;font-size:10px;margin:2px;color:#c9d1d9}}
.btn{{display:inline-block;padding:5px 12px;border-radius:6px;font-size:11px;border:1px solid #30363d;background:#21262d;color:#c9d1d9;text-decoration:none;margin:3px 0}}
.btn:hover{{background:#30363d}}
</style>
</head>
<body>
<div id="hdr">
  <h1>🧠 {title}</h1>
  <div class="st">
    <span>{len(nodes)} nodes &nbsp;·&nbsp; {len(edges)} connections &nbsp;·&nbsp; {n_comm} clusters</span>
    <span class="badge EX">AST: {stats['extracted']}</span>
    <span class="badge IN">Inferred: {stats['inferred']}</span>
    <span class="badge AM">Ambiguous: {stats['ambiguous']}</span>
  </div>
</div>
<div id="body">
  <div id="sb">
    <input type="text" id="search" placeholder="Search nodes...">
    <h4>Topics</h4>{topic_filters}
    <h4>Clusters</h4>{comm_filters}
    <h4>Navigate</h4>
    <a href="../index.html" class="btn">🏠 Dashboard</a>
    <a href="report.md" class="btn">📊 Report</a>
    <a href="../wiki/index.md" class="btn">📖 Wiki</a>
    <a href="../tiers/index.md" class="btn">📦 Tiers</a>
  </div>
  <div id="graph"></div>
  <div id="info">
    <h3 id="i-title"></h3>
    <div class="meta" id="i-meta"></div>
    <div class="kws" id="i-kws"></div>
    <div id="i-links"></div>
  </div>
</div>
<script>
const ND={vis_nodes};
const ED={vis_edges};
const nodes=new vis.DataSet(ND);
const edges=new vis.DataSet(ED);
const net=new vis.Network(document.getElementById('graph'),{{nodes,edges}},{{
  physics:{{stabilization:{{iterations:200}},barnesHut:{{gravitationalConstant:-4500,springLength:140,springConstant:0.03}}}},
  interaction:{{hover:true,tooltipDelay:300}},
  nodes:{{shape:'dot',borderWidth:1.5}},
  edges:{{smooth:{{type:'continuous'}}}},
}});
net.on('click',p=>{{
  if(!p.nodes.length){{document.getElementById('info').style.display='none';return;}}
  const id=p.nodes[0];
  const n=ND.find(x=>x.id===id);
  document.getElementById('info').style.display='block';
  document.getElementById('i-title').textContent=n.label;
  const m=n.title;
  document.getElementById('i-meta').innerHTML=
    `Topic: <b>${{n.topic}}</b><br>`+
    `File: <code style="font-size:10px">${{m.match(/File: ([^<]+)/)?.[1]||''}}</code><br>`+
    `Connections: <b>${{m.match(/Connections: ([^<]+)/)?.[1]||''}}</b> &nbsp; Cluster: ${{n.community}}`;
  const kws=(m.split('Keywords: ')[1]||'').split('<')[0];
  document.getElementById('i-kws').innerHTML=kws.split(', ').map(k=>`<span>${{k}}</span>`).join('');
  document.getElementById('i-links').innerHTML=
    `<a href="../wiki/${{n.topic}}/${{id}}.md" class="btn">📖 Wiki page</a>`+
    `<a href="../tiers/entity/${{id}}.md" class="btn">📦 Deep context</a>`;
}});
document.getElementById('search').addEventListener('input',function(){{
  const q=this.value.toLowerCase();
  nodes.forEach(n=>nodes.update({{id:n.id,hidden:!!q&&!n.label.toLowerCase().includes(q)}}));
}});
function applyFilters(){{
  const at=[...document.querySelectorAll('.ft:checked')].map(c=>c.value);
  const ac=[...document.querySelectorAll('.fc:checked')].map(c=>+c.value);
  nodes.forEach(n=>nodes.update({{id:n.id,hidden:!at.includes(n.topic)||!ac.includes(n.community)}}));
}}
document.querySelectorAll('.ft,.fc').forEach(cb=>cb.addEventListener('change',applyFilters));
</script>
</body>
</html>"""
    (graph_dir / 'graph.html').write_text(html, encoding='utf-8')
    print(f'  Graph: graph/graph.html ({len(nodes)} nodes, {len(edges)} edges, {n_comm} clusters)')
    return stats


# ── Dashboard ─────────────────────────────────────────────────────────────────

def build_dashboard(entities: List[dict], edges: List[dict], out_dir: Path, title: str, stats: dict):
    topics: Dict[str, list] = defaultdict(list)
    for e in entities:
        topics[e['topic']].append(e)
    cards = ''.join(
        f'<div class="card"><div style="width:10px;height:10px;border-radius:50%;background:{TOPIC_COLORS.get(t,"#888")};flex-shrink:0"></div>'
        f'<div><div class="ct">{t.replace("-"," ").title()}</div><div class="cs">{len(ents)} entities</div></div></div>'
        for t, ents in sorted(topics.items()))
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><title>{title}</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;background:#0d1117;color:#e6edf3;min-height:100vh}}
.hero{{padding:60px 40px 30px;text-align:center}}
.hero h1{{font-size:32px;font-weight:700;color:#fff;margin-bottom:8px}}
.hero p{{font-size:13px;color:#8b949e;margin-bottom:4px}}
.stats{{display:flex;gap:14px;justify-content:center;margin:26px 0;flex-wrap:wrap}}
.stat{{background:#161b22;border:1px solid #30363d;border-radius:10px;padding:16px 26px;text-align:center}}
.sn{{font-size:28px;font-weight:700;color:#388bfd}}.sl{{font-size:12px;color:#8b949e;margin-top:2px}}
.conf{{display:flex;gap:10px;justify-content:center;margin-bottom:24px;flex-wrap:wrap}}
.badge{{padding:4px 12px;border-radius:10px;font-size:12px;font-weight:600}}
.EX{{background:#1f6feb33;color:#58a6ff}}.IN{{background:#388bfd22;color:#79c0ff}}.AM{{background:#f0883e22;color:#f0883e}}
.actions{{display:flex;gap:12px;justify-content:center;margin-bottom:40px;flex-wrap:wrap}}
.btn{{padding:10px 22px;border-radius:8px;font-size:14px;font-weight:600;text-decoration:none;transition:opacity .2s}}
.bp{{background:#388bfd;color:#fff}}.bs{{background:#21262d;color:#e6edf3;border:1px solid #30363d}}
.btn:hover{{opacity:.85}}
.sec{{max-width:860px;margin:0 auto;padding:0 40px 50px}}
.sec h2{{font-size:15px;font-weight:600;margin-bottom:12px;color:#c9d1d9}}
.cards{{display:grid;grid-template-columns:repeat(auto-fill,minmax(175px,1fr));gap:10px}}
.card{{background:#161b22;border:1px solid #30363d;border-radius:8px;padding:12px 14px;display:flex;align-items:center;gap:10px}}
.ct{{font-size:13px;font-weight:500}}.cs{{font-size:11px;color:#8b949e;margin-top:2px}}
footer{{text-align:center;padding:20px;font-size:11px;color:#484f58}}
</style>
</head>
<body>
<div class="hero">
  <h1>🧠 {title}</h1>
  <p>Second Brain T v1.0 &nbsp;·&nbsp; Second Brain T v1.0</p>
  <p style="color:#484f58;font-size:11px">{datetime.now().strftime('%Y-%m-%d %H:%M')}</p>
  <div class="stats">
    <div class="stat"><div class="sn">{len(entities)}</div><div class="sl">Entities</div></div>
    <div class="stat"><div class="sn">{len(edges)}</div><div class="sl">Connections</div></div>
    <div class="stat"><div class="sn">{len(topics)}</div><div class="sl">Topics</div></div>
    <div class="stat"><div class="sn">{stats.get("communities",0)}</div><div class="sl">Clusters</div></div>
  </div>
  <div class="conf">
    <span class="badge EX">AST/Imports: {stats.get("extracted",0)}</span>
    <span class="badge IN">Inferred: {stats.get("inferred",0)}</span>
    <span class="badge AM">Ambiguous: {stats.get("ambiguous",0)}</span>
  </div>
  <div class="actions">
    <a href="graph/graph.html" class="btn bp">🔗 Knowledge Graph</a>
    <a href="wiki/index.md" class="btn bs">📖 Obsidian Wiki</a>
    <a href="tiers/index.md" class="btn bs">📦 Context Tiers</a>
    <a href="graph/report.md" class="btn bs">📊 Report</a>
  </div>
</div>
<div class="sec"><h2>Topics</h2><div class="cards">{cards}</div></div>
<footer>Second Brain T v1.0 &nbsp;·&nbsp; Second Brain T v1.0 &nbsp;·&nbsp; MIT</footer>
</body>
</html>"""
    (out_dir / 'index.html').write_text(html, encoding='utf-8')
    print(f'  Dashboard: index.html')


# ── Claude skill installer ─────────────────────────────────────────────────────

def install_skill(target: Path):
    skill_dir = target / '.claude' / 'skills' / 'second-brain-t'
    skill_dir.mkdir(parents=True, exist_ok=True)
    (skill_dir / 'SKILL.md').write_text(
        '---\nname: second-brain-t\ndescription: Use when the user wants to build a knowledge base, '
        'analyze a project, create a wiki, or build a knowledge graph.\n---\n\n'
        '# /sbt\n\nRun: `python3 build.py <folder> [--pdf] [--no-cache]`\n\n'
        '## Query\n1. Load `output/tiers/index.md` (Tier 0)\n'
        '2. Load `output/tiers/topic/{topic}.md` (Tier 1)\n'
        '3. Load `output/tiers/entity/{id}.md` (Tier 2)\n'
        '4. See `output/graph/report.md` for connections\n',
        encoding='utf-8')
    cmd_dir = target / '.claude' / 'commands'
    cmd_dir.mkdir(parents=True, exist_ok=True)
    (cmd_dir / 'sbt.md').write_text(
        '# /sbt — Query Second Brain T\n\n'
        '```\n/sbt           # Stats\n/sbt compile   # Rebuild\n'
        '/sbt search X  # Search tiers\n/sbt gaps      # Find orphans\n```\n\n'
        'Always start from `output/tiers/index.md`, then drill into topic or entity files.\n',
        encoding='utf-8')


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    target, title, use_cache = parse_args()
    if not target:
        print('Usage: python3 build.py <folder> [--title "Name"] [--no-cache]')
        sys.exit(1)
    if not target.exists():
        print(f'Error: {target} not found')
        sys.exit(1)

    OUT = Path(__file__).parent / 'output'
    OUT.mkdir(parents=True, exist_ok=True)

    print(f'\nSecond Brain T v1.0')
    print(f'  Source : {target}')
    print(f'  Output : {OUT}\n')

    # Always auto-convert PDFs, Word docs, and PowerPoint files
    pdfs = collect_pdfs(target)
    docx_files = collect_docx(target)
    pptx_files = collect_pptx(target)
    if pdfs or docx_files or pptx_files:
        total = len(pdfs) + len(docx_files) + len(pptx_files)
        print(f'Found {total} document(s) to convert (PDF:{len(pdfs)} DOCX:{len(docx_files)} PPTX:{len(pptx_files)})...')
        n = convert_pdfs(target) + convert_docx(target) + convert_pptx(target)
        print(f'  {n} converted\n')

    files = collect_files(target)
    print(f'Scanning... {len(files)} files found')

    cache = Cache(OUT, enabled=use_cache)
    cached = sum(1 for f in files if cache.load(f))
    print(f'  Cache: {cached}/{len(files)} files already cached\n')

    print('Extracting entities...')
    entities = build_entities(files, target, cache)
    code_count = sum(1 for e in entities if e.get('is_code'))
    print(f'  {len(entities)} entities ({code_count} with AST parsing)\n')

    print('Finding relationships...')
    edges = find_relationships(entities)
    ex = sum(1 for e in edges if e['type'] == 'EXTRACTED')
    inf = sum(1 for e in edges if e['type'] == 'INFERRED')
    amb = sum(1 for e in edges if e['type'] == 'AMBIGUOUS')
    print(f'  {len(edges)} connections — AST: {ex} | Inferred: {inf} | Ambiguous: {amb}\n')

    print('Detecting communities...')
    communities = detect_communities(entities, edges)
    n_comm = max(communities.values(), default=0) + 1
    print(f'  {n_comm} clusters\n')

    print('Building outputs...')
    build_tiers(entities, edges, OUT / 'tiers', title)
    build_wiki(entities, edges, OUT / 'wiki', title)
    stats = build_graph(entities, edges, OUT / 'graph', title, communities)
    stats['communities'] = n_comm
    build_dashboard(entities, edges, OUT, title, stats)
    install_skill(target)

    (OUT / 'freshness.json').write_text(json.dumps({
        'compiled_at': datetime.now().isoformat(),
        'entities': len(entities), 'edges': len(edges),
        'tool': 'Second Brain T v1.0',
    }, indent=2), encoding='utf-8')

    print(f'\n✓ Done!\n')
    print(f'  Open  : {OUT / "index.html"}')
    print(f'  Wiki  : Open {OUT / "wiki"} in Obsidian')
    print(f'  Query : Ask Claude using /sbt\n')


if __name__ == '__main__':
    main()
